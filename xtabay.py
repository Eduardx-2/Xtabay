from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.slide import Slide
import json
from pptx.dml.color import RGBColor


setConfig = r"conf/text_box.json"



class ConfImagesMode:
    def __init__(self):
        self.text_x = ""
        #---image
        self.image_x = ""
        self.image_y = ""
        self.image_w = ""
        self.image_h = ""

    def _text_r(self):
        return self.text_x

    #---image
    def _image_x(self):
        return Inches(self.image_x)

    def _image_y(self):
        return Inches(self.image_y)

    def _image_w(self):
        return Inches(self.image_w)

    def _image_h(self):
        return Inches(self.image_h)
    
    def _mode_coordenadas(self,clave:str,valor:str):
        com_mode = {
            "tree-d": {
                "DCS":{ #imagen a la derecha
                    "texto": [0.8, 1.2, 6.0, 4.8],
                    "imagen": [7.2, 1.2, 5.3, 4.8]
                },
                "CT": {
                    "texto": [0.5, 1.2, 4.5, 4.5],
                    "imagen": [5.3, 1.2, 4.0, 4.0]
                }
            },
            "tree-i": {
                "DCS":{
                    "texto": [6.6, 1.2, 5.8, 4.8],
                    "imagen": [0.8, 1.2, 5.3, 4.8]
                },
                "CT":{#imagen a la izquierda 
                    "texto": [5.0, 1.2, 4.2, 4.5],
                    "imagen": [0.6, 1.2, 3.8, 4.0]
                }
            },
            "image-a":{ #imagen arriba
                "DCS":{"texto": [0.8, 4.5, 11.7, 2.0],
                "imagen": [1.5, 0.8, 10.0, 3.2]},
                "CT":{
                   "texto": [0.7, 5.0, 8.5, 1.2],
                   "imagen": [1.2, 0.8, 7.5, 3.5]
                }
            },
            "image-l":{#Imagen grande con descripción lateral
                "DCS":{"texto": [9.3, 1.0, 3.0, 5.0],
                "imagen": [0.5, 1.0, 8.2, 5.5]},
                "CT":{
                    "texto": [7.2, 1.2, 2.0, 4.5],
                    "imagen": [0.3, 0.8, 6.3, 5.5]
                }
            }
            
        }
        try:
            self.text_x = com_mode[clave][valor]["texto"]
        except KeyError:
            return False
    
        self.image_x,self.image_y,self.image_w,self.image_h = com_mode[clave][valor]["imagen"]
        return True

class Xtabayxc:
    def __init__(self):
        self.presentacion = Presentation()
        self.conf_ = Configurations()
        self.config_lecter = self.conf_.inicie_config()
        self.colors = self.conf_._color_set()
       

    def imagen_pptx_text_data(self,mode:str, slide:Slide,image,valor):
        config_pptx_f = ConfImagesMode() #Si tengo mas metodos, este match se puede expander tontamente
        #Tengo que corregir este match, no es necesario saber el modo no? para que 
        mode_seg = config_pptx_f._mode_coordenadas(mode,valor)
        if mode_seg:
            slide.shapes.add_picture(image,config_pptx_f._image_x(),
                                        config_pptx_f._image_y(),
                                        width=config_pptx_f._image_w(),
                                        height=config_pptx_f._image_h())
       
        #UN SOLO METODO PARA DISTINTAS FORMAS DE IMAGEN CON TEXTO
      
        return [slide,config_pptx_f._text_r(),config_pptx_f._text_r()]

    def subtitle_cong(self,texto,slide_pgh:Slide,conf:list,conf_lec,color,typeMod:str):#Posicion 2:4
        config = self.config_lecter[conf_lec[0]][conf_lec[1]]["subtitle"]
        textoInch = slide_pgh.shapes.add_textbox(Inches(config[0]),Inches(config[1]),Inches(config[2]),Inches(config[3]))
        conft_text = textoInch.text_frame
        pf_text = conft_text.paragraphs[0]
        if "anidado" in typeMod.lower():
            if isinstance(texto,(list)):
                text_anidado = " | ".join(texto)

                pf_text.text = text_anidado
                pf_text.font.color.rgb = color
                pf_text.font.size = Pt(int(conf[1][0]))
                pf_text.space_after = Pt(4)
            else:
                print("[!] ERROR LIST SUBTITLE MODE")
        else:
           
            pf_text.text = texto
            match conf[0]:
                case "bold": pf_text.font.bold = True
                case "italic": pf_text.font.italic = True
                case _ : print("[-] ERROR SUBTTITLE")
    
            pf_text.font.color.rgb = color
            pf_text.font.size = Pt(int(conf[1][0]))
        

    def title_presentation_conf(self,texto_data:str,confg:list, slide_pag:Slide,title_mod:list,color) -> Slide: #recibe la comnfiguración -> title recibe las coordenadas del titulo
        #OK lo obtengo pero solo tengo la clvae del titulo no la clave principal
        caseConf = self.config_lecter[title_mod[0]][title_mod[1]]["title"]
        textBox = slide_pag.shapes.add_textbox(Inches(caseConf[0]),Inches(caseConf[1]),Inches(caseConf[2]),Inches(caseConf[3]))
        texto = textBox.text_frame
        texto.word_wrap = True
        paragh = texto.paragraphs[0]
        paragh.text = texto_data
        match confg[0]:
            case "bold": paragh.font.bold  = True
            case "italic": paragh.font.italic = True
            case _ : print("[-] ERROR")
        
        paragh.font.color.rgb = color
        paragh.font.size = Pt(int(confg[1][0]))
        return slide_pag

    def recurs_body_pptx(self,slide:Slide,mode_cont:list,text:str,envConfig:list): 
        config_normal_text = self.config_lecter[mode_cont[0]][mode_cont[1]]["content"]
        text_BoxBody = slide.shapes.add_textbox(Inches(config_normal_text[0]),
                                                        Inches(config_normal_text[1]),
                                                        Inches(config_normal_text[2]),
                                                        Inches(config_normal_text[3]))
        textContent = text_BoxBody.text_frame
        textContent.word_wrap = True
        graphContent = textContent.paragraphs[0]
        graphContent.text = text
        graphContent.font.color.rgb = envConfig[1]
        graphContent.font.size = Pt(envConfig[0])
        

    def texto_body_pptx_normal(self,body_text,confg:list,slide:Slide,mode_cont:list,color,conf_d:list,mode):
        textBoxBody = slide.shapes.add_textbox(Inches(conf_d[0]),Inches(conf_d[1]),Inches(conf_d[2]),Inches(conf_d[3]))
        texto_body = textBoxBody.text_frame #ESTE TIENE QUE IR EN GRAPH
        texto_body.word_wrap = True
        body_graph = texto_body.paragraphs[0]   
        
        if "graph" in mode:
            
            if isinstance(body_text,(list)):
                for gph in body_text:
                    textgph = texto_body.add_paragraph()
                    textgph.text = f"° {gph}"
                    textgph.font.size = Pt(int(confg[1][0]))
                    textgph.font.color.rgb = color
                    textgph.space_after = Pt(8)
            else:
                print("[-] ERROR LIST")
        elif "i_normal" in mode:
           
            self.recurs_body_pptx(slide,
                                  mode_cont,
                                  body_text,
                                  [int(confg[1][0]),color])
        else:

            body_graph.text = body_text
            match confg[0]:
                case "bold": body_graph.font.bold = True
                case "italic": body_graph.font.italic = True
                case _ : print("[-] ERROR")
            body_graph.font.color.rgb = color
            body_graph.font.size = Pt(int(confg[1][0]))
        return slide


    def inicie_data_pptx(self,confText,name_file):
        configsession = confText #QUE HAGO AQUI? TENGO QUE LEER PARAM_JSON 
        modpapel_ = configsession["informacion"]["space-pag"] 
        colors_met = configsession["colors_letra"]
        listaValor = [colors_met["titulo"],colors_met["subtitulo"],colors_met["body"]] #COLORES
        colors_data = [self.colors[d] for d in listaValor]
        pags = configsession['informacion']["pags"]
        self.conf_.format_config(modpapel_)
        if self.conf_.ESTADO:
            self.presentacion.slide_height = Inches(int(self.conf_.space_height_conf())) 
            self.presentacion.slide_width = Inches(int(self.conf_.space_width_conf()))
    
        for i in range(1,pags+1):

            titles_pptx = configsession["informacion"]["title"][f"title_{i}"]
            letra_titles_ = [titles_pptx["typeLetra"],titles_pptx["tituloFont"]]
    
            pagf_session = configsession["informacion"]["text_body"][f"pag_{i}"] #OBTIENE LAS KEYS DE PAG
            new_confg_lect = [pagf_session["typeLetra"],pagf_session["cuerpoContent"]]
            if pagf_session["verifyImage"].lower() == "true":
                image_data = self.imagen_pptx_text_data(pagf_session["mode"],
                                                        self.presentacion.slides.add_slide(self.presentacion.slide_layouts[6]),
                                                        pagf_session["image"],modpapel_)
                
                self.texto_body_pptx_normal(pagf_session["content"],new_confg_lect,image_data[0],
                                            [modpapel_,configsession["informacion"]["space-mod"]],
                                            colors_data[2],image_data[1],pagf_session["format-text"])
                
            elif pagf_session["verifyImage"].lower() == "false":
                text2 = self.texto_body_pptx_normal(pagf_session["content"],new_confg_lect,
                                                self.presentacion.slides.add_slide(self.presentacion.slide_layouts[6]),
                                                [modpapel_,configsession["informacion"]["space-mod"]],colors_data[2],
                                                [0.7, 0.5, 8.6, 1.0],pagf_session["format-text"])
           
            if titles_pptx["content"] is not None:
                   
            
                text_final_body = self.title_presentation_conf(titles_pptx["content"],letra_titles_,
                                            text2,
                                            [modpapel_,configsession["informacion"]["space-mod"]],
                                            colors_data[0])

            if titles_pptx["subtitle"] is not None:
                
                self.subtitle_cong(titles_pptx["subtitle"]["content"],text_final_body,
                                   [titles_pptx["subtitle"]["typeLetra"],titles_pptx["subtitle"]["font"]],
                                   [modpapel_,configsession["informacion"]["space-mod"]],
                                   colors_data[2],
                                   titles_pptx["subtitle"]["type"])
            

        self.presentacion.save(f"{name_file}.pptx")
       
class Configurations:
    def __init__(self):
        self.WIDTH = ""
        self.HEIGHT = ""
        self.ESTADO = False
        

    def space_width_conf(self):
        return self.WIDTH

    def space_height_conf(self):
        return self.HEIGHT

    def estado_get(self):
        return self.ESTADO
    
    def inicie_config(self):
        with open(setConfig,'r') as popen_of:
            json_set = json.load(popen_of)
            return json_set

    
            
    def _color_set(self) -> dict:
        colors = {
            "negro": RGBColor(0, 0, 0),
            "blanco": RGBColor(255, 255, 255),

            "gris_carbon": RGBColor(51, 51, 51),
            "gris_oscuro": RGBColor(68, 68, 68),
            "gris_grafito": RGBColor(85, 85, 85),
            "gris_medio": RGBColor(102, 102, 102),

            "azul_marino": RGBColor(15, 23, 42),
            "azul_real": RGBColor(37, 99, 235),
            "azul_cielo": RGBColor(96, 165, 250),

            "verde_esmeralda": RGBColor(22, 163, 74),
            "verde_bosque": RGBColor(22, 101, 52),

            "violeta": RGBColor(109, 40, 217),
            "purpura": RGBColor(147, 51, 234),

            "rojo": RGBColor(220, 38, 38),
            "naranja": RGBColor(234, 88, 12),

            "dorado": RGBColor(202, 138, 4)
        }
        return colors
    
    def format_config(self,format_:str):
        slide_ = {
            "CT": [10.0, 7.5],
            "DCS": [13.333, 7.5],
            "DCDS": [10.0, 6.25],
            "A4H": [11.69, 8.27],
            "A4V": [8.27, 11.69],
            "LEH": [11.0, 8.5],
            "LEV": [8.5, 11.0]
        }
        search = slide_[format_]
        self.WIDTH = search[0]
        self.HEIGHT = search[1]
        self.ESTADO = True

